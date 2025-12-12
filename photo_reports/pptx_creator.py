# photo_reports/pptx_creator.py
import ast
import calendar
import logging
import re
from collections import defaultdict
from datetime import datetime
from pathlib import Path
from typing import Dict, Tuple, List, Any, Sequence

import pandas as pd
from PIL import ImageFont, Image, ImageDraw
from pptx import Presentation
from pptx.util import Pt, Cm, Inches
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT, MSO_VERTICAL_ANCHOR

from configs.config import (
    TITLE_CONTENT,
    logger,
    FOLDER_PATH,
    ENABLED_WORK_TYPE,
    START_DATE,
    END_DATE,
    REPORT_MODE,
    PHOTO_LABELS_ENABLED,
)
from configs.env_utils import get_env_json, get_env_str
from configs.pptx_config import (
    ADDRESS_PREFIX,
    CONSTRUCTION_MULTI_PREFIX,
    CONSTRUCTION_NUMBER_TEMPLATE,
    CONSTRUCTION_SINGLE_TEMPLATE,
    DATE_RANGE_SEPARATOR,
    LABEL_FONT_PATH,
    LABEL_FONT_SIZE,
    LABEL_FONT_FAMILY,
    MONTH_NAMES,
    PHOTO_HEIGHT,
    LABEL_BOX,
    LABEL_VERTICAL_SHIFT,
    LABEL_PHOTO_GAP,
    LABEL_HEIGHT,
    REPORT_FALLBACK_NAME,
    WORK_STAGE_FALLBACK_TITLES,
    WORK_TITLES,
    PHOTO_TOP,
    SIDE_MARGIN,
    GAP,
    CONTENT_WIDTH_RATIO,
    MAX_GAP_MULTIPLIER,
    WORK_SLIDE_SIDE_MARGIN,
    WORK_SLIDE_TOP_MARGIN,
    WORK_SLIDE_BOTTOM_MARGIN,
    WORK_SLIDE_MIN_GAP,
)
from configs.utils_config import RU_MONTH, STAGE_ORDER
from configs.work_profiles import WORK_ACTIVITY_MAP as WORK_ACTIVITY_DEFINITIONS
from utils.utils import ensure_pptx_compatible

_font_cache: Dict[int, ImageFont.FreeTypeFont] = {}
_image_size_cache: Dict[str, Tuple[int, int]] = {}
_dummy_img = Image.new("RGB", (1, 1))
_dummy_draw = ImageDraw.Draw(_dummy_img)


def _compute_spacing(widths: List[int], available_width: int, base_gap: int, max_gap_multiplier: float) -> tuple[int, int]:
    """Return (gap, free_space) for a row of photos given desired width budget."""
    if not widths:
        return base_gap, available_width

    slots = max(len(widths) - 1, 0)
    total_width_no_gap = sum(widths)

    if slots == 0:
        free_space = max(0, available_width - total_width_no_gap)
        return 0, free_space

    gap = base_gap
    max_gap = int(base_gap * max_gap_multiplier)

    min_total = total_width_no_gap + gap * slots
    extra_space = max(0, available_width - min_total)

    if max_gap > gap and slots:
        allowable_increase = min(max_gap - gap, extra_space // slots)
        gap += allowable_increase
        extra_space -= allowable_increase * slots

    total_width = total_width_no_gap + gap * slots
    free_space = max(0, available_width - total_width)
    return gap, free_space


def _config_date_to_datetime(value: Any) -> datetime:
    if isinstance(value, datetime):
        return value
    if isinstance(value, str):
        return datetime.strptime(value, "%d.%m.%Y")
    raise TypeError(f"Unsupported date type: {type(value)!r}")


def _month_label(dt: datetime) -> str:
    return f"{dt.year}-{dt.month:02d}"


def _month_bounds(dt: datetime) -> Tuple[datetime, datetime]:
    last_day = calendar.monthrange(dt.year, dt.month)[1]
    month_start = datetime(dt.year, dt.month, 1)
    month_end = datetime(dt.year, dt.month, last_day)
    return month_start, month_end



REPORT_ACTIVITY_MAP: Dict[str, Tuple[str, str]] = {
    "daily": (
        get_env_str("DAILY_ACTIVITY_NAME", "Control Inspection"),
        get_env_str("DAILY_ACTIVITY_SHORT", "CI"),
    ),
}

WORK_ACTIVITY_MAP: Dict[str, Tuple[str, str]] = WORK_ACTIVITY_DEFINITIONS


def _extract_address() -> str:
    prefix = ADDRESS_PREFIX
    for line in TITLE_CONTENT.splitlines():
        if line.strip().startswith(prefix):
            addr = line.split(":", 1)[1].strip()
            for junk in ("\u00ab", "\u00bb", '"'):
                addr = addr.replace(junk, "")
            addr = addr.replace("\u00A0", " ")
            return addr
    return REPORT_FALLBACK_NAME

REPORT_ADDRESS = _extract_address()


def _sanitize_filename_fragment(value: str) -> str:
    invalid_chars = '<>:"/\\|?*'
    cleaned = "".join("_" if ch in invalid_chars else ch for ch in value)
    cleaned = cleaned.strip().rstrip('.')
    return cleaned or "report"


def _resolve_activity(report_mode: str, work_type: str) -> Tuple[str, str]:
    if report_mode == "daily":
        return REPORT_ACTIVITY_MAP["daily"]
    normalized = (work_type or "").strip()
    if not normalized:
        return ("Work Activities", "Work")
    return WORK_ACTIVITY_MAP.get(normalized, (normalized, normalized))


def _month_display_name(dt: datetime) -> str:
    return MONTH_NAMES[dt.month] if 1 <= dt.month <= 12 else f"{dt.month:02d}"


def get_cached_image_size(path: str) -> Tuple[int, int]:
    size = _image_size_cache.get(path)
    if size is None:
        with Image.open(path) as img:
            size = (img.width, img.height)
        _image_size_cache[path] = size
    return size


def get_aspect_ratio(path: str) -> float:
    w, h = get_cached_image_size(path)
    return w / h


def shrink_text_to_fit(
        text: str,
        width_pt: float,
        height_pt: float,
        font_path: str,
        max_font_size: int,
        min_font_size: int = 8
) -> int:
    low, high = min_font_size, max_font_size
    best = min_font_size

    while low <= high:
        mid = (low + high) // 2
        if mid not in _font_cache:
            try:
                _font_cache[mid] = ImageFont.truetype(font_path, mid)
            except IOError:
                _font_cache[mid] = ImageFont.load_default()
        font = _font_cache[mid]
        bbox = _dummy_draw.textbbox((0, 0), text, font=font)
        text_w = bbox[2] - bbox[0]
        text_h = bbox[3] - bbox[1]

        if text_w <= width_pt and text_h <= height_pt:
            best = mid
            low = mid + 1
        else:
            high = mid - 1

    return best


def add_label(slide, left, top, width, height, text) -> None:
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE

    width_pt = width / 12700
    height_pt = height / 12700
    font_size = shrink_text_to_fit(text, width_pt, height_pt, LABEL_FONT_PATH, LABEL_FONT_SIZE)

    p = tf.paragraphs[0]
    p.text = text
    p.font.name = LABEL_FONT_FAMILY
    p.font.size = Pt(font_size)
    p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER


def is_landscape(path: str) -> bool:
    w, h = get_cached_image_size(path)
    return w > h


def add_blank_slide(prs: Presentation):
    blank_layout = None
    for layout in prs.slide_layouts:
        try:
            placeholders = layout.placeholders
        except AttributeError:
            placeholders = ()
        if len(placeholders) == 0:
            blank_layout = layout
            break

    if blank_layout is None:
        blank_layout = prs.slide_layouts[5]

    slide = prs.slides.add_slide(blank_layout)
    for shape in list(slide.shapes):
        if getattr(shape, "is_placeholder", False):
            element = shape._element
            parent = element.getparent()
            if parent is not None:
                parent.remove(element)
    return slide


def format_construction_line(numbers: Sequence[int]) -> str:
    """Return header text for one or multiple constructions."""
    unique: List[int] = []
    seen = set()
    for num in numbers:
        if num in seen:
            continue
        seen.add(num)
        unique.append(num)

    if not unique:
        return ""
    if len(unique) == 1:
        return CONSTRUCTION_SINGLE_TEMPLATE.format(num=unique[0])
    numbers = ", ".join(CONSTRUCTION_NUMBER_TEMPLATE.format(num=n) for n in unique)
    return f"{CONSTRUCTION_MULTI_PREFIX}{numbers}"


def format_period_text(start: datetime, end: datetime) -> str:
    """Return a well-spaced period string using the configured separator."""
    separator = DATE_RANGE_SEPARATOR.strip() or "-"
    text = f"{start.strftime('%d.%m.%Y')} {separator} {end.strftime('%d.%m.%Y')}"
    return " ".join(text.split())


def add_title_slide(prs: Presentation, activity_text: str, period_text: str) -> None:
    slide = add_blank_slide(prs)
    tb = slide.shapes.add_textbox(0, 0, prs.slide_width, prs.slide_height)
    tf = tb.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE

    p = tf.paragraphs[0]
    p.text = TITLE_CONTENT.format(activity=activity_text, period=period_text)
    p.font.name = LABEL_FONT_FAMILY
    p.font.size = Pt(LABEL_FONT_SIZE)
    p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER


def add_photo_with_label(slide, photo_path: str, left: int, photo_top: int, label_text: str) -> None:
    aspect = get_aspect_ratio(photo_path)
    pic_h = PHOTO_HEIGHT
    pic_w = int(pic_h * aspect)
    slide.shapes.add_picture(photo_path, left, photo_top, width=pic_w, height=pic_h)

    if PHOTO_LABELS_ENABLED:
        label_w, label_h = map(Cm, LABEL_BOX)
        label_left = left + (pic_w - label_w) // 2
        label_top = photo_top - label_h - int(LABEL_VERTICAL_SHIFT + LABEL_PHOTO_GAP)
        add_label(slide, label_left, label_top, label_w, label_h, label_text)


def insert_slides_adaptive(prs: Presentation, items: List[Tuple[str, str]]) -> None:
    """Render one horizontal photo row per slide at a fixed height (PHOTO_HEIGHT).
    Side margins (`SIDE_MARGIN`) stay symmetric and the gap between photos is never smaller than `GAP`.
    When the next image no longer fits in the available width, a new slide is created.
    """
    if not items:
        return

    side_margin = int(SIDE_MARGIN)
    min_gap = int(GAP)
    photo_top = int(PHOTO_TOP)

    # Use the full slide width minus the side margins (CONTENT_WIDTH_RATIO is ignored).
    usable_width = int(prs.slide_width) - 2 * side_margin
    if usable_width <= 0:
        usable_width = 1
    content_left = side_margin

    def draw_row(slide, row_items: List[Tuple[str, str, int]]) -> None:
        if not row_items:
            return
        gaps_total = min_gap * (len(row_items) - 1) if len(row_items) > 1 else 0
        total = sum(w for _, _, w in row_items) + gaps_total
        start_x = content_left + max(0, (usable_width - total) // 2)

        x = start_x
        for i, (path, label, w) in enumerate(row_items):
            add_photo_with_label(slide, path, int(x), photo_top, label)
            if i < len(row_items) - 1:
                x += w + min_gap

    slide = add_blank_slide(prs)
    row: List[Tuple[str, str, int]] = []
    used = 0

    for path, label in items:
        w = int(PHOTO_HEIGHT * get_aspect_ratio(path))
        # Guard against ultra-wide photos so a single image does not dominate the row.
        w = min(w, usable_width)

        extra = w if not row else (min_gap + w)
        if row and used + extra > usable_width:
            draw_row(slide, row)
            slide = add_blank_slide(prs)
            row = [(path, label, w)]
            used = w
        else:
            if not row:
                row = [(path, label, w)]
                used = w
            else:
                row.append((path, label, w))
                used += min_gap + w

    if row:
        draw_row(slide, row)
    return
    """Portrait photos occupy one row per slide at the same fixed height (PHOTO_HEIGHT).
    Margins stay equal, the minimum spacing is `GAP`, and the row remains horizontally centred.
    If the next photo overflows the remaining width, a new slide is started.
    """
    if not items:
        return

    side_margin = int(SIDE_MARGIN)
    min_gap = int(GAP)
    photo_top = int(PHOTO_TOP)

    full_available = int(prs.slide_width) - 2 * side_margin
    content_width = max(1, int(round(full_available * float(CONTENT_WIDTH_RATIO))))
    content_left = side_margin + max(0, (full_available - content_width) // 2)
    usable_width = content_width

    def draw_row(slide, row_items: List[Tuple[str, str, int]]) -> None:
        if not row_items:
            return
        gaps_total = min_gap * (len(row_items) - 1) if len(row_items) > 1 else 0
        total = sum(w for _, _, w in row_items) + gaps_total
        start_x = content_left + max(0, (usable_width - total) // 2)
        x = start_x
        for i, (path, label, w) in enumerate(row_items):
            add_photo_with_label(slide, path, int(x), photo_top, label)
            if i < len(row_items) - 1:
                x += w + min_gap

    slide = add_blank_slide(prs)
    row: List[Tuple[str, str, int]] = []
    used = 0

    for path, label in items:
        w = int(PHOTO_HEIGHT * get_aspect_ratio(path))
        w = min(w, usable_width)
        extra = w if not row else (min_gap + w)

        if row and used + extra > usable_width:
            draw_row(slide, row)
            slide = add_blank_slide(prs)
            row = [(path, label, w)]
            used = w
        else:
            if not row:
                row = [(path, label, w)]
                used = w
            else:
                row.append((path, label, w))
                used += min_gap + w

    if row:
        draw_row(slide, row)
    return
    """Portrait photos occupy one row per slide at the same fixed height (PHOTO_HEIGHT).
    Margins stay equal, the minimum spacing is `GAP`, and the row remains horizontally centred.
    If the next photo overflows the remaining width, a new slide is started.
    """
    if not items:
        return

    side_margin = int(SIDE_MARGIN)
    min_gap = int(GAP)
    photo_top = int(PHOTO_TOP)

    full_available = int(prs.slide_width) - 2 * side_margin
    content_width = max(1, int(round(full_available * float(CONTENT_WIDTH_RATIO))))
    content_left = side_margin + max(0, (full_available - content_width) // 2)
    usable_width = content_width

    def draw_row(slide, row_items: List[Tuple[str, str, int]]) -> None:
        if not row_items:
            return
        gaps_total = min_gap * (len(row_items) - 1) if len(row_items) > 1 else 0
        total = sum(w for _, _, w in row_items) + gaps_total
        start_x = content_left + max(0, (usable_width - total) // 2)
        x = start_x
        for i, (path, label, w) in enumerate(row_items):
            add_photo_with_label(slide, path, int(x), photo_top, label)
            if i < len(row_items) - 1:
                x += w + min_gap

    slide = add_blank_slide(prs)
    row: List[Tuple[str, str, int]] = []
    used = 0

    for path, label in items:
        w = int(PHOTO_HEIGHT * get_aspect_ratio(path))
        w = min(w, usable_width)  # Clamp width if the photo exceeds the usable area.
        extra = w if not row else (min_gap + w)

        if row and used + extra > usable_width:
            draw_row(slide, row)
            slide = add_blank_slide(prs)
            row = [(path, label, w)]
            used = w
        else:
            if not row:
                row = [(path, label, w)]
                used = w
            else:
                row.append((path, label, w))
                used += min_gap + w

    if row:
        draw_row(slide, row)
    """Place selected photos in a single row per slide.
    Keep the height fixed (PHOTO_HEIGHT), preserve equal margins, and apply at least `GAP` spacing.
    If the next photo cannot fit with the required gap, render the current row and start a new slide.
    """
    if not items:
        return

    side_margin = int(SIDE_MARGIN)
    min_gap = int(GAP)
    photo_top = int(PHOTO_TOP)

    full_available = int(prs.slide_width) - 2 * side_margin
    content_width = max(1, int(round(full_available * float(CONTENT_WIDTH_RATIO))))
    content_left = side_margin + max(0, (full_available - content_width) // 2)
    usable_width = content_width

    def draw_row(slide, row_items: List[Tuple[str, str, int]]) -> None:
        if not row_items:
            return
        n = len(row_items)
        gaps_total = min_gap * (n - 1) if n > 1 else 0
        total = sum(w for _, _, w in row_items) + gaps_total
        start_x = content_left + max(0, (usable_width - total) // 2)

        x = start_x
        for i, (path, label, w) in enumerate(row_items):
            # Keep width tied to the fixed height: width = height * aspect ratio.
            add_photo_with_label(slide, path, int(x), photo_top, label)
            if i < n - 1:
                x += w + min_gap

    # Collect photos along with their calculated widths.
    slide = add_blank_slide(prs)
    row: List[Tuple[str, str, int]] = []
    used = 0  # Track row width already occupied by photos and gaps.

    for path, label in items:
        w = int(PHOTO_HEIGHT * get_aspect_ratio(path))  # Width for the fixed-height layout.
        extra = w if not row else (min_gap + w)
        if row and used + extra > usable_width:
            # Flush the current row and start a new slide.
            draw_row(slide, row)
            slide = add_blank_slide(prs)
            row = [(path, label, w)]
            used = w
        else:
            if not row:
                row = [(path, label, w)]
                used = w
            else:
                row.append((path, label, w))
                used += min_gap + w

    if row:
        draw_row(slide, row)
    if not items:
        return

    side_margin = int(SIDE_MARGIN)
    min_gap = int(GAP)
    photo_top = int(PHOTO_TOP)

    full_available = int(prs.slide_width) - 2 * side_margin
    content_width = max(1, int(round(full_available * float(CONTENT_WIDTH_RATIO))))
    content_left = side_margin + max(0, (full_available - content_width) // 2)

    usable_width = content_width

    slide = add_blank_slide(prs)
    row: List[Tuple[str, str, int]] = []
    row_width = 0

    def flush_row(curr_slide, row_items: List[Tuple[str, str, int]]):
        if not row_items:
            return curr_slide
        n = len(row_items)
        gaps_total = min_gap * (n - 1) if n > 1 else 0
        total = sum(w for _, _, w in row_items) + gaps_total
        start_x = content_left + max(0, (usable_width - total) // 2)

        x = start_x
        for i, (p, lbl, w) in enumerate(row_items):
            add_photo_with_label(curr_slide, p, int(x), photo_top, lbl)
            if i < n - 1:
                x += w + min_gap

        # When the row is full, create a new blank slide.
        return add_blank_slide(prs)

    for path, label in items:
        pic_w = int(PHOTO_HEIGHT * get_aspect_ratio(path))
        pic_w = min(pic_w, usable_width)  # Clamp overly wide photos to the usable width.

        needed = pic_w if not row else (row_width + min_gap + pic_w)
        if row and needed > usable_width:
            slide = flush_row(slide, row)
            row = [(path, label, pic_w)]
            row_width = pic_w
        else:
            if not row:
                row = [(path, label, pic_w)]
                row_width = pic_w
            else:
                row.append((path, label, pic_w))
                row_width += min_gap + pic_w

    if row:
        n = len(row)
        gaps_total = min_gap * (n - 1) if n > 1 else 0
        total = sum(w for _, _, w in row) + gaps_total
        start_x = content_left + max(0, (usable_width - total) // 2)

        x = start_x
        for i, (p, lbl, w) in enumerate(row):
            add_photo_with_label(slide, p, int(x), photo_top, lbl)
            if i < n - 1:
                x += w + min_gap
    """Lay out photos in one row per slide with a fixed height (PHOTO_HEIGHT).
    Use equal side margins (`SIDE_MARGIN`) and enforce a minimum gap of `GAP` between photos.
    When the next image exceeds the available width (including the gap), draw the row and move to a new slide.
    Rows never stretch beyond the centred content area.
    """
    if not items:
        return

    left_margin = int(SIDE_MARGIN)
    right_margin = int(SIDE_MARGIN)
    min_gap = int(GAP)
    photo_top = int(PHOTO_TOP)

    usable_width = int(prs.slide_width) - left_margin - right_margin
    usable_width = max(1, usable_width)

    slide = add_blank_slide(prs)
    row: List[Tuple[str, str, int]] = []  # (path, label, width_int)
    row_width = 0

    def flush_row(curr_slide, row_items: List[Tuple[str, str, int]]):
        if not row_items:
            return curr_slide
        n = len(row_items)
        gaps_total = min_gap * (n - 1) if n > 1 else 0
        total = sum(w for _, _, w in row_items) + gaps_total
        start_x = left_margin + max(0, (usable_width - total) // 2)

        x = start_x
        for i, (p, lbl, w) in enumerate(row_items):
            add_photo_with_label(curr_slide, p, int(x), photo_top, lbl)
            if i < n - 1:
                x += w + min_gap

        # Flush the current row and return a new blank slide.
        return add_blank_slide(prs)

    for path, label in items:
        pic_w = int(PHOTO_HEIGHT * get_aspect_ratio(path))
        if pic_w > usable_width:
            pic_w = usable_width  # Clamp the width to the usable area.

        needed = pic_w if not row else (row_width + min_gap + pic_w)
        if row and needed > usable_width:
            slide = flush_row(slide, row)
            row = [(path, label, pic_w)]
            row_width = pic_w
        else:
            if not row:
                row = [(path, label, pic_w)]
                row_width = pic_w
            else:
                row.append((path, label, pic_w))
                row_width += min_gap + pic_w

    # Centre the final row horizontally before drawing it.
    if row:
        n = len(row)
        gaps_total = min_gap * (n - 1) if n > 1 else 0
        total = sum(w for _, _, w in row) + gaps_total
        start_x = left_margin + max(0, (usable_width - total) // 2)

        x = start_x
        for i, (p, lbl, w) in enumerate(row):
            add_photo_with_label(slide, p, int(x), photo_top, lbl)
            if i < n - 1:
                x += w + min_gap
    """Arrange photos in rows: height fixed to PHOTO_HEIGHT and margins (`SIDE_MARGIN`) stay symmetrical.
    Each photo claims its width plus the gap; if it would overflow, we move the row to a new slide.
    The row is centred horizontally after placement.
    """
    if not items:
        return

    left_margin = int(SIDE_MARGIN)
    right_margin = int(SIDE_MARGIN)  # Mirror the left margin.
    min_gap = int(GAP)
    photo_top = int(PHOTO_TOP)

    usable_width = int(prs.slide_width) - left_margin - right_margin
    if usable_width <= 0:
        usable_width = 1

    slide = add_blank_slide(prs)
    x = left_margin
    first_in_row = True

    for path, label in items:
        # Compute width for a fixed-height photo.
        pic_w = int(PHOTO_HEIGHT * get_aspect_ratio(path))

        # Determine how much horizontal space this photo requires.
        needed = pic_w if first_in_row else (min_gap + pic_w)

        # Start a new slide when adding the photo would overflow the usable width.
        if (x - left_margin) + needed > usable_width:
            slide = add_blank_slide(prs)
            x = left_margin
            first_in_row = True
            needed = pic_w

        # Position the photo within the row while respecting the gap.
        left = x if first_in_row else x + min_gap
        add_photo_with_label(slide, path, int(left), photo_top, label)

        # Advance the cursor for the next photo.
        x = left + pic_w
        first_in_row = False


def insert_slide(prs: Presentation, items: List[Tuple[str, str]], orientation: str = 'vertical'):
    slide = add_blank_slide(prs)
    side_margin = int(SIDE_MARGIN)
    photo_top = int(PHOTO_TOP)

    if orientation == 'vertical':
        pic_widths = [int(PHOTO_HEIGHT * get_aspect_ratio(path)) for path, _ in items]
        count = len(items)
        if count == 0:
            return

        full_available = prs.slide_width - 2 * side_margin
        usable_width = int(full_available * CONTENT_WIDTH_RATIO)
        usable_width = max(usable_width, 1)
        left_offset = side_margin + (full_available - usable_width) // 2
        base_gap = int(GAP)

        if count > 1:
            column_width = usable_width // count
            if column_width > 0 and max(pic_widths) <= column_width:
                for idx, (path, label) in enumerate(items):
                    width = pic_widths[idx]
                    column_left = left_offset + idx * column_width
                    inner_offset = (column_width - width) // 2
                    add_photo_with_label(slide, path, int(column_left + inner_offset), photo_top, label)
                return

        gap_between, free_space = _compute_spacing(pic_widths, usable_width, base_gap, MAX_GAP_MULTIPLIER)
        x = left_offset + free_space // 2

        for idx, (path, label) in enumerate(items):
            left = int(x)
            add_photo_with_label(slide, path, left, photo_top, label)
            if idx < count - 1:
                x += pic_widths[idx] + gap_between
            else:
                x += pic_widths[idx]

    else:
        for path, label in items:
            aspect = get_aspect_ratio(path)
            adjusted_width = int(PHOTO_HEIGHT * aspect)
            max_width = prs.slide_width - 2 * side_margin
            if adjusted_width > max_width:
                adjusted_width = max_width
            photo_left = int((prs.slide_width - adjusted_width) / 2)
            slide.shapes.add_picture(path, photo_left, photo_top, width=adjusted_width, height=PHOTO_HEIGHT)
            if PHOTO_LABELS_ENABLED:
                label_h = int(LABEL_HEIGHT)
                label_top = photo_top - label_h - int(LABEL_VERTICAL_SHIFT + LABEL_PHOTO_GAP)
                add_label(slide, photo_left, label_top, adjusted_width, label_h, label)



def _render_daily_month(
    df_month: pd.DataFrame,
    out_path: Path,
    month_label: str,
    activity_text: str,
    period_text: str,
) -> None:
    session_order = {"day": 0, "night": 1, "single": 2}
    df_sorted = df_month.copy()
    if "session" not in df_sorted.columns:
        df_sorted["session"] = ""
    df_sorted["__session_key"] = (
        df_sorted["session"].fillna("single").astype(str).str.lower().map(session_order).fillna(len(session_order))
    )
    df_sorted.sort_values(["date_dt", "__session_key", "cn0"], inplace=True)
    df_sorted.drop(columns=["__session_key"], inplace=True)

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    add_title_slide(prs, activity_text, period_text)

    buffer: List[Tuple[str, str]] = []
    prev_date: str | None = None
    prev_session: str | None = None

    for _, row in df_sorted.iterrows():
        curr_date = row["base_date"]
        curr_session = str(row.get("session", "") or "single").strip().lower()
        if prev_date is not None and (curr_date != prev_date or curr_session != prev_session):
            if buffer:
                insert_slides_adaptive(prs, buffer)
                buffer.clear()
        prev_date = curr_date
        prev_session = curr_session

        orig = Path(row["path"])
        base, ext, folder = orig.stem, orig.suffix, orig.parent

        final_path = None
        for suffix in ("_filtered_stamped", "_stamped", "_filtered"):
            cand = folder / f"{base}{suffix}{ext}"
            if cand.exists():
                final_path = str(cand)
                break
        if final_path is None:
            final_path = str(orig)
        if not Path(final_path).exists():
            logger.warning("File not found for slide: %s", final_path)
            continue

        date_line = row["base_date"]
        nums = row["construction_number"]
        num_line = format_construction_line(nums)
        work_line = WORK_TITLES.get(row.get("work_type", ""), "")
        header_lines = [date_line]
        if num_line:
            header_lines.append(num_line)
        if work_line:
            header_lines.append(work_line)
        header = "\n".join(header_lines)

        if is_landscape(final_path):
            if buffer:
                insert_slides_adaptive(prs, buffer)
                buffer.clear()
            insert_slide(prs, [(final_path, header)], orientation="horizontal")
        else:
            buffer.append((final_path, header))

    # Flush any remaining photos into a final slide.
    if buffer:
        insert_slides_adaptive(prs, buffer)

    prs.save(out_path)
    logger.info("Saved daily presentation for %s to %s", month_label, out_path)


def create_daily_presentation(report_mode: str | None = None, work_type: str | None = None) -> None:
    excel_path = FOLDER_PATH / "parsed_records.xlsx"
    df = pd.read_excel(excel_path, dtype=str)

    df["construction_number"] = df["construction_number"].apply(
        lambda x: ast.literal_eval(x) if isinstance(x, str) else []
    )

    df_daily = df[df["stage"].isna() | (df["stage"] == "detected")].copy()
    if df_daily.empty:
        logger.warning("No detected-stage records found for daily reports.")
        return

    effective_mode = report_mode or REPORT_MODE
    effective_work_type = work_type if work_type is not None else ENABLED_WORK_TYPE
    activity_text, file_label = _resolve_activity(effective_mode, effective_work_type)
    address_text = REPORT_ADDRESS

    df_daily["date_dt"] = pd.to_datetime(df_daily["base_date"], format="%d.%m.%Y", errors="coerce")
    df_daily = df_daily.dropna(subset=["date_dt"])
    df_daily["cn0"] = df_daily["construction_number"].apply(lambda numbers: numbers[0] if numbers else 0)

    start_dt = _config_date_to_datetime(START_DATE)
    end_dt = _config_date_to_datetime(END_DATE)
    mask = (df_daily["date_dt"] >= start_dt) & (df_daily["date_dt"] <= end_dt)
    df_daily = df_daily.loc[mask]

    if df_daily.empty:
        logger.warning("No daily records between %s and %s", START_DATE, END_DATE)
        return

    grouped = df_daily.groupby(df_daily["date_dt"].dt.to_period("M"), sort=True)
    for period, df_month in grouped:
        if df_month.empty:
            continue
        month_dt = period.to_timestamp()
        month_label = _month_label(month_dt)
        month_start, month_end = _month_bounds(month_dt)
        period_start = max(month_start, start_dt)
        period_end = min(month_end, end_dt)
        period_text = format_period_text(period_start, period_end)
        month_name = _month_display_name(month_dt)
        sanitized_address = _sanitize_filename_fragment(address_text)
        out_name = f"{sanitized_address}. {file_label}. {month_name}.pptx"
        out_path = Path(FOLDER_PATH) / out_name
        _render_daily_month(df_month, out_path, month_label, activity_text, period_text)


def _render_work_month(
    records: List[Dict[str, Any]],
    month_label: str,
    out_path: Path,
    activity_text: str,
    period_text: str,
) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    add_title_slide(prs, activity_text, period_text)

    grouped: Dict[str, Dict[Tuple[int, ...], Dict[str, Dict[str, List[Dict[str, Any]]]]]] = defaultdict(
        lambda: defaultdict(lambda: defaultdict(lambda: defaultdict(list)))
    )
    for r in records:
        grouped[r["base_date"]][r["cn_numbers"]][r["stage"]][r["stage_title"]].append(r)

    def stage_sort_key(stage: str) -> Tuple[int, str]:
        return STAGE_ORDER.get(stage, 99), stage or ""

    for date in sorted(grouped, key=lambda d: datetime.strptime(d, "%d.%m.%Y")):
        cn_map = grouped[date]
        for cn_key in sorted(cn_map):
            stage_map = cn_map[cn_key]
            for stage in sorted(stage_map, key=stage_sort_key):
                title_groups = stage_map[stage]
                for stage_title in sorted(title_groups.keys()):
                    photos = title_groups[stage_title]
                    if not photos:
                        continue

                    left_margin = int(WORK_SLIDE_SIDE_MARGIN)
                    right_margin = int(WORK_SLIDE_SIDE_MARGIN)
                    header_offset = int(LABEL_VERTICAL_SHIFT + LABEL_PHOTO_GAP)
                    header_top = max(0, int(WORK_SLIDE_TOP_MARGIN) - int(LABEL_HEIGHT) - header_offset)
                    usable_width = int(prs.slide_width) - left_margin - right_margin
                    usable_width = max(1, usable_width)
                    usable_width_float = float(usable_width)
                    photos_area_top = int(WORK_SLIDE_TOP_MARGIN)
                    photos_area_bottom = int(prs.slide_height) - int(WORK_SLIDE_BOTTOM_MARGIN)
                    available_height = max(1, photos_area_bottom - photos_area_top)
                    available_height_float = float(available_height)
                    base_gap = max(1, int(WORK_SLIDE_MIN_GAP))
                    min_gap = float(base_gap)
                    target_row_height = min(float(int(PHOTO_HEIGHT)), available_height_float)

                    header_lines = [date]
                    cn_line = format_construction_line(cn_key)
                    if cn_line:
                        header_lines.append(cn_line)
                    header_lines.append(stage_title)
                    header_text = "\n".join(header_lines)

                    photo_infos: List[Tuple[str, float]] = []
                    for rec in photos:
                        orig = Path(rec["path"])
                        base = orig.stem
                        ext = orig.suffix
                        folder = orig.parent

                        final = None
                        for suf in ("_filtered_stamped", "_stamped", "_filtered"):
                            for ext2 in (ext, ".png", ".jpg", ".jpeg"):
                                cand = folder / f"{base}{suf}{ext2}"
                                if cand.exists():
                                    final = cand
                                    break
                            if final:
                                break
                        if final is None:
                            final = orig

                        final_path = ensure_pptx_compatible(str(final))
                        ratio = get_aspect_ratio(final_path)
                        photo_infos.append((final_path, ratio))

                    idx = 0
                    while idx < len(photo_infos):
                        slide = add_blank_slide(prs)
                        if PHOTO_LABELS_ENABLED:
                            add_label(slide, left_margin, header_top, usable_width, LABEL_HEIGHT, header_text)

                        # Collect images for the current row until no more fit
                        row_items: List[Tuple[str, float]] = []
                        row_width = 0.0
                        while idx < len(photo_infos):
                            path, ratio = photo_infos[idx]
                            width = ratio * float(PHOTO_HEIGHT)
                            if width > usable_width_float:
                                width = usable_width_float
                            needed = width if not row_items else width + min_gap
                            if row_items and row_width + needed > usable_width_float:
                                break
                            row_items.append((path, width))
                            row_width += needed
                            idx += 1

                        if not row_items:
                            path, ratio = photo_infos[idx]
                            width = min(ratio * float(PHOTO_HEIGHT), usable_width_float)
                            row_items.append((path, width))
                            idx += 1

                        # Scale row to fit into the content width
                        widths_float = [width for _, width in row_items]
                        gap_int = base_gap if len(row_items) > 1 else 0
                        gap_total = gap_int * max(0, len(row_items) - 1)

                        usable_for_photos = max(1.0, usable_width_float - gap_total)
                        sum_widths = sum(widths_float)
                        scale = 1.0
                        if sum_widths > 0 and usable_for_photos < sum_widths:
                            scale = usable_for_photos / sum_widths

                        row_height = max(1, int(round(float(PHOTO_HEIGHT) * scale)))
                        widths_int = [max(1, int(round(width * scale))) for width in widths_float]
                        gap_int = max(1, gap_int) if len(row_items) > 1 else 0

                        total_drawn_width = sum(widths_int) + gap_int * (len(row_items) - 1)
                        if total_drawn_width > usable_width:
                            overflow = total_drawn_width - usable_width
                            widths_int[-1] = max(1, widths_int[-1] - overflow)
                            total_drawn_width = sum(widths_int) + gap_int * (len(row_items) - 1)

                        content_left = left_margin
                        content_right = left_margin + usable_width
                        start_x = content_left + max(0, (usable_width - total_drawn_width) // 2)
                        y = photos_area_top

                        x = start_x
                        placed_shapes: List[Any] = []
                        for index_item, (path, width_int) in enumerate(zip((p for p, _ in row_items), widths_int)):
                            shape = slide.shapes.add_picture(
                                path,
                                int(round(x)),
                                int(round(y)),
                                width=width_int,
                                height=row_height,
                            )
                            placed_shapes.append(shape)
                            x += width_int
                            if gap_int and index_item < len(widths_int) - 1:
                                x += gap_int

                        if placed_shapes:
                            actual_left = min(shape.left for shape in placed_shapes)
                            actual_right = max(shape.left + shape.width for shape in placed_shapes)
                            actual_width = max(1, actual_right - actual_left)
                            ideal_left = content_left + max(0, (usable_width - actual_width) // 2)
                            delta = ideal_left - actual_left
                            if delta:
                                for shape in placed_shapes:
                                    shape.left = int(round(shape.left + delta))

                            updated_left = min(shape.left for shape in placed_shapes)
                            updated_right = max(shape.left + shape.width for shape in placed_shapes)
                            if updated_left < content_left:
                                shift = content_left - updated_left
                                for shape in placed_shapes:
                                    shape.left = int(shape.left + shift)
                                updated_right += shift
                            if updated_right > content_right:
                                shift = content_right - updated_right
                                for shape in placed_shapes:
                                    shape.left = int(shape.left + shift)

    prs.save(out_path)
    logger.info("Saved work presentation for %s (%s) to %s", month_label, activity_text, out_path)


def create_work_presentation(data: List[Dict[str, Any]], work_type: str | None = None) -> None:
    """Generate per-work monthly presentations."""
    allowed_stages = {"detected", "in_progress", "fixed"}
    fallback_titles = WORK_STAGE_FALLBACK_TITLES

    start_dt = _config_date_to_datetime(START_DATE)
    end_dt = _config_date_to_datetime(END_DATE)

    selected_work_type = (work_type or ENABLED_WORK_TYPE or "").strip()
    activity_text, file_label = _resolve_activity("work", selected_work_type)
    address_text = REPORT_ADDRESS

    buckets: Dict[str, List[Dict[str, Any]]] = defaultdict(list)
    month_periods: Dict[str, Tuple[datetime, datetime]] = {}

    for record in data:
        if selected_work_type and record.get("work_type") != selected_work_type:
            continue

        stage = str(record.get("stage") or "").strip().lower()
        if stage not in allowed_stages:
            continue

        cns = record.get("construction_number")
        if isinstance(cns, str):
            try:
                cns = ast.literal_eval(cns)
            except Exception:
                cns = []
        if not isinstance(cns, list):
            cns = []

        cn_numbers: List[int] = []
        seen = set()
        for raw in cns:
            token = str(raw)
            if not token.isdigit():
                continue
            value = int(token)
            if value in seen:
                continue
            seen.add(value)
            cn_numbers.append(value)
        if not cn_numbers:
            continue
        cn_key = tuple(cn_numbers)

        base_date = str(record.get("base_date") or record.get("date") or "").strip()
        if not base_date:
            continue
        try:
            base_dt = datetime.strptime(base_date, "%d.%m.%Y")
        except ValueError:
            continue
        if not (start_dt <= base_dt <= end_dt):
            continue

        stage_folder = Path(record["path"]).parent.name
        stage_title = re.sub(r"^\d+\s+", "", stage_folder).strip() or fallback_titles.get(stage, stage)

        payload = {
            "base_date": base_date,
            "cn_numbers": cn_key,
            "stage": stage,
            "stage_title": stage_title,
            "path": record["path"],
        }
        month_key = _month_label(base_dt)
        buckets[month_key].append(payload)

        month_start, month_end = _month_bounds(base_dt)
        period_start = max(month_start, start_dt)
        period_end = min(month_end, end_dt)
        prev_bounds = month_periods.get(month_key)
        if prev_bounds:
            prev_start, prev_end = prev_bounds
            period_start = min(prev_start, period_start)
            period_end = max(prev_end, period_end)
        month_periods[month_key] = (period_start, period_end)

    if not buckets:
        logger.warning(
            "No work records for %s between %s and %s",
            selected_work_type or "(any)",
            START_DATE,
            END_DATE,
        )
        return

    for month_key in sorted(buckets.keys()):
        month_records = buckets[month_key]
        if not month_records:
            continue
        period_start, period_end = month_periods.get(month_key, (start_dt, end_dt))
        period_text = format_period_text(period_start, period_end)
        month_dt = datetime.strptime(month_key, "%Y-%m")
        month_name = _month_display_name(month_dt)
        sanitized_address = _sanitize_filename_fragment(address_text)
        out_name = f"{sanitized_address}. {file_label}. {month_name}.pptx"
        out_path = Path(FOLDER_PATH) / out_name
        _render_work_month(month_records, month_key, out_path, activity_text, period_text)
